import os
import sys
import code
import numpy as np
import subprocess

class Commit(object):
	def __init__(self, author, date, commit_message):
		self.author = author
		self.date = date
		self.commit_message = commit_message

def getCommits(path):
	truePath = os.path.dirname(os.path.abspath(__file__))
	#code.interact(local=locals())
	os.chdir(path)
	s = 'git whatchanged --since="30 days ago" > ' + truePath + '/' + 'commits.txt'
	subprocess.call(s, shell=True)
	os.chdir(truePath)



def buildCommits():

	with open("commits.txt", "r") as f:
	    content = f.readlines()

	lines = [line.rstrip('\n') for line in content]

	commit = ""
	author = ""
	message = ""
	date = ""
	first = 0
	allCommits = []

	for each in lines:
		each = each.strip()
		if len(each) == 0:
			continue
		if each.split()[0] == 'commit':
			if first != 0:
				com = Commit(author, date, message)
				allCommits.append(com)
				commit = ""
				author = ""
				message = ""
				date = ""
			else:
				commit = each
				first = 1
		elif each.split()[0] == 'Author:':
			author = each
		elif each.split()[0] == 'Date:':
			date = each
		elif each[0] == ':':
			continue
		else:
			message += each + ' '
	return allCommits

def filterCommits(allCommits, authorName):
	filteredCommits = []
	for each in allCommits:
		if authorName in each.author:
			filteredCommits.append(each)
	return filteredCommits

def printCommits(allCommits):
	for each in allCommits:
		print each.author
		print each.date
		print each.commit_message
		print " "
		print " "

def writeCommit(commits):
	f = open("filteredCommits.txt", 'w')
	for each in commits:
		f.write(each.author + '\r\n')
		f.write(each.date + '\r\n')
		f.write(each.commit_message + '\r\n')
		f.write('\r\n')
		f.write('\r\n')


def buildPPTX(commits):
	from pptx import Presentation
	prs = Presentation()
	title_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	subtitle = slide.placeholders[1]

	title.text = "Monthly Presentation"
	subtitle.text = "Using git2ppt!"

	prs.save('test.pptx')
	commits.reverse()
	for ii,each in enumerate(commits):
		print ii
		bullet_slide_layout = prs.slide_layouts[5]
		slide = prs.slides.add_slide(bullet_slide_layout)
		shapes = slide.shapes
		title = slide.shapes.title
		title.text = each.commit_message


	prs.save('test.pptx')

def main():
	print "git2ppt -- Written by Michael Morehead Nov. 2016"
	print "usage: python git2ppt.py [path/to/repo] [makePPT = 0]"
	print " set makePPT to 1 to try to make a ppt"
	if len(sys.argv) < 2:
		sys.exit(0)
	else:
		makePPT = 0
		if len(sys.argv) == 3:
			makePPT = sys.argv[2]
		pathToRepo = sys.argv[1]
		getCommits(pathToRepo)
		allCommits = buildCommits()
		filtered = filterCommits(allCommits, "mmorehea")

		printCommits(filtered)
		writeCommit(filtered)

		if makePPT:
			buildPPTX(filtered)


if __name__ == "__main__":
    main()
